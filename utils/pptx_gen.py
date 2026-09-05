from pptx import Presentation
from pptx.util import Inches
from db import supabase
from datetime import datetime

def generate_analysis_deck() -> str:
    """Generates a PPTX analysis deck of recent sales."""
    try:
        prs = Presentation()
        
        # Slide 1: Title
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Kirana Store Business Analysis"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        
        # Fetch data for analysis (e.g., all finalized bills)
        res_bills = supabase.table("bills").select("*").eq("status", "finalized").execute()
        total_sales = sum(float(b['total_amount']) for b in res_bills.data)
        total_tax = sum(float(b['total_tax']) for b in res_bills.data)
        
        # Slide 2: Summary
        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide2.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "High-Level Sales Summary"
        tf = body_shape.text_frame
        tf.text = f"Total Orders: {len(res_bills.data)}"
        p = tf.add_paragraph()
        p.text = f"Gross Sales Volume: ₹{total_sales:.2f}"
        p = tf.add_paragraph()
        p.text = f"Total GST Collected: ₹{total_tax:.2f}"
        
        pptx_path = "sales_analysis_deck.pptx"
        prs.save(pptx_path)
        return pptx_path
        
    except Exception as e:
        return f"Error generating PPTX: {str(e)}"

from langchain_core.tools import tool
@tool
def generate_analysis_deck_tool() -> str:
    """Generates a PPTX analysis deck and returns the file path."""
    return generate_analysis_deck()
